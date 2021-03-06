# encoding: utf-8

"""Test suite for pptx.slides module."""

from __future__ import absolute_import

import pytest

from lxml import objectify
from mock import ANY, call, MagicMock, Mock

from pptx.opc import package
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.packuri import PackURI
from pptx.opc.package import Part, _Relationship
from pptx.oxml.ns import namespaces
from pptx.oxml.presentation import CT_SlideId, CT_SlideIdList
from pptx.parts.slides import (
    _BaseSlide, Slide, SlideCollection, SlideLayout, SlideMaster
)
from pptx.presentation import Package, Presentation
from pptx.shapes.shapetree import ShapeCollection

from ..unitutil import (
    absjoin, class_mock, instance_mock, loose_mock, method_mock,
    parse_xml_file, serialize_xml, test_file_dir
)


test_image_path = absjoin(test_file_dir, 'python-icon.jpeg')
test_pptx_path = absjoin(test_file_dir, 'test.pptx')

nsmap = namespaces('a', 'r', 'p')


def actual_xml(elm):
    objectify.deannotate(elm, cleanup_namespaces=True)
    return serialize_xml(elm, pretty_print=True)


def _sldLayout1():
    path = absjoin(test_file_dir, 'slideLayout1.xml')
    sldLayout = parse_xml_file(path).getroot()
    return sldLayout


def _sldLayout1_shapes():
    sldLayout = _sldLayout1()
    spTree = sldLayout.xpath('./p:cSld/p:spTree', namespaces=nsmap)[0]
    shapes = ShapeCollection(spTree)
    return shapes


class Describe_BaseSlide(object):

    def it_knows_the_name_of_the_slide(self, base_slide):
        """_BaseSlide.name value is correct"""
        # setup ------------------------
        base_slide._element = _sldLayout1()
        # exercise ---------------------
        name = base_slide.name
        # verify -----------------------
        expected = 'Title Slide'
        actual = name
        msg = "expected '%s', got '%s'" % (expected, actual)
        assert actual == expected, msg

    def it_provides_access_to_the_shapes_on_the_slide(self):
        """_BaseSlide.shapes is expected size after _load()"""
        # setup ------------------------
        path = absjoin(test_file_dir, 'slide1.xml')
        with open(path, 'r') as f:
            blob = f.read()
        base_slide = _BaseSlide.load(None, None, blob, None)
        # exercise ---------------------
        shapes = base_slide.shapes
        # verify -----------------------
        assert len(shapes) == 9

    def it_can_add_an_image_part_to_the_slide(self, base_slide_fixture):
        # fixture ----------------------
        base_slide, img_file_, image_, rId_ = base_slide_fixture
        # exercise ---------------------
        image, rId = base_slide._add_image(img_file_)
        # verify -----------------------
        base_slide._package._images.add_image.assert_called_once_with(
            img_file_)
        base_slide.relate_to.assert_called_once_with(image, RT.IMAGE)
        assert image is image_
        assert rId is rId_

    def it_knows_it_is_the_part_its_child_objects_belong_to(
            self, base_slide):
        assert base_slide.part is base_slide

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def base_slide(self):
        partname = PackURI('/foo/bar.xml')
        return _BaseSlide(partname, None, None, None)

    @pytest.fixture
    def base_slide_fixture(self, request, base_slide):
        # mock _BaseSlide._package._images.add_image() train wreck
        img_file_ = loose_mock(request, name='img_file_')
        image_ = loose_mock(request, name='image_')
        pkg_ = loose_mock(request, name='_package', spec=Package)
        pkg_._images.add_image.return_value = image_
        base_slide._package = pkg_
        # mock _BaseSlide.relate_to()
        rId_ = loose_mock(request, name='rId_')
        method_mock(request, _BaseSlide, 'relate_to', return_value=rId_)
        return base_slide, img_file_, image_, rId_


class DescribeSlide(object):

    def it_establishes_a_relationship_to_its_slide_layout_on_construction(
            self, relate_to_):
        """Slide(slidelayout) adds relationship slide->slidelayout"""
        # setup ------------------------
        slidelayout = SlideLayout(None, None, _sldLayout1(), None)
        partname = PackURI('/ppt/slides/slide1.xml')
        # exercise ---------------------
        slide = Slide.new(slidelayout, partname, None)
        # verify ----------------------
        slide.relate_to.assert_called_once_with(slidelayout, RT.SLIDE_LAYOUT)

    # def it_creates_a_minimal_sld_element_on_construction(self, slide):
    #     """Slide._element is minimal sld on construction"""
    #     # setup ------------------------
    #     slidelayout = SlideLayout(None, None, _sldLayout1())
    #     partname = PackURI('/ppt/slides/slide1.xml')
    #     slide = Slide.new(slidelayout, partname)
    #     path = absjoin(test_file_dir, 'minimal_slide.xml')
    #     # exercise ---------------------
    #     elm = slide._element
    #     # verify -----------------------
    #     with open(path, 'r') as f:
    #         expected_xml = f.read()
    #     assert actual_xml(elm) == expected_xml

    # def it_has_slidelayout_property_of_none_on_construction(self, slide):
    #     """Slide.slidelayout property None on construction"""
    #     assert slide.slidelayout is None

    # def it_sets_slidelayout_on_load(self, slide):
    #     """Slide._load() sets slidelayout"""
    #     # setup ------------------------
    #     path = absjoin(test_file_dir, 'slide1.xml')
    #     slidelayout = Mock(name='slideLayout')
    #     slidelayout.partname = '/ppt/slideLayouts/slideLayout1.xml'
    #     rel = Mock(name='pptx.package.Relationship')
    #     rel.rId = 'rId1'
    #     rel.reltype = RT.SLIDE_LAYOUT
    #     rel.target = slidelayout
    #     pkgpart = Mock(name='pptx.package.Part')
    #     with open(path, 'rb') as f:
    #         pkgpart.blob = f.read()
    #     pkgpart.relationships = [rel]
    #     part_dict = {slidelayout.partname: slidelayout}
    #     slide_ = slide.load(pkgpart, part_dict)
    #     # exercise ---------------------
    #     retval = slide_.slidelayout
    #     # verify -----------------------
    #     expected = slidelayout
    #     actual = retval
    #     msg = "expected: %s, got %s" % (expected, actual)
    #     assert actual == expected, msg

    def it_knows_the_minimal_element_xml_for_a_slide(self, slide):
        """Slide._minimal_element generates correct XML"""
        # setup ------------------------
        path = absjoin(test_file_dir, 'minimal_slide.xml')
        # exercise ---------------------
        sld = slide._minimal_element()
        # verify -----------------------
        with open(path, 'r') as f:
            expected_xml = f.read()
        assert actual_xml(sld) == expected_xml

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def relate_to_(self, request):
        return method_mock(request, Part, 'relate_to')

    @pytest.fixture
    def slide(self):
        return Slide(None, None, None, None)


class DescribeSlideCollection(object):

    def it_supports_indexed_access(self, slides_with_slide_parts_, rIds_):
        slides, slide_, slide_2_ = slides_with_slide_parts_
        rId_, rId_2_ = rIds_
        # verify -----------------------
        assert slides[0] is slide_
        assert slides[1] is slide_2_
        slides._sldIdLst.__getitem__.assert_has_calls(
            [call(0), call(1)]
        )
        slides._prs.related_parts.__getitem__.assert_has_calls(
            [call(rId_), call(rId_2_)]
        )

    def it_raises_on_slide_index_out_of_range(self, slides):
        with pytest.raises(IndexError):
            slides[2]

    def it_can_iterate_over_the_slides(self, slides, slide_, slide_2_):
        assert [s for s in slides] == [slide_, slide_2_]

    def it_supports_len(self, slides):
        assert len(slides) == 2

    def it_can_add_a_new_slide(self, slides, slidelayout_, Slide_, slide_):
        slide = slides.add_slide(slidelayout_)
        Slide_.new.assert_called_once_with(
            slidelayout_, PackURI('/ppt/slides/slide3.xml'),
            slides._prs.package
        )
        slides._prs.relate_to.assert_called_once_with(slide_, RT.SLIDE)
        slides._sldIdLst.add_sldId.assert_called_once_with(ANY)
        assert slide is slide_

    def it_knows_the_next_available_slide_partname(
            self, slides_with_slide_parts_):
        slides = slides_with_slide_parts_[0]
        expected_partname = PackURI('/ppt/slides/slide3.xml')
        partname = slides._next_partname
        assert isinstance(partname, PackURI)
        assert partname == expected_partname

    def it_can_assign_partnames_to_the_slides(
            self, slides, slide_, slide_2_):
        slides.rename_slides()
        assert slide_.partname == '/ppt/slides/slide1.xml'
        assert slide_2_.partname == '/ppt/slides/slide2.xml'

    # fixtures -------------------------------------------------------
    #
    #   slides
    #   |
    #   +- ._sldIdLst = [sldId_, sldId_2_]
    #   |                |       |
    #   |                |       +- .rId = rId_2_
    #   |                |
    #   |                +- .rId = rId_
    #   +- ._prs
    #       |
    #       +- .related_parts = {rId_: slide_, rId_2_: slide_2_}
    #
    # ----------------------------------------------------------------

    @pytest.fixture
    def prs_(self, request, rel_, related_parts_):
        prs_ = instance_mock(request, Presentation)
        prs_.load_rel.return_value = rel_
        prs_.related_parts = related_parts_
        return prs_

    @pytest.fixture
    def rel_(self, request, rId_):
        return instance_mock(request, _Relationship, rId=rId_)

    @pytest.fixture
    def related_parts_(self, request, rIds_, slide_parts_):
        """
        Return pass-thru mock dict that both operates as a dict an records
        calls to __getitem__ for call asserts.
        """
        rId_, rId_2_ = rIds_
        slide_, slide_2_ = slide_parts_
        slide_rId_map = {rId_: slide_, rId_2_: slide_2_}

        def getitem(key):
            return slide_rId_map[key]

        related_parts_ = MagicMock()
        related_parts_.__getitem__.side_effect = getitem
        return related_parts_

    @pytest.fixture
    def rename_slides_(self, request):
        return method_mock(request, SlideCollection, 'rename_slides')

    @pytest.fixture
    def rId_(self, request):
        return 'rId1'

    @pytest.fixture
    def rId_2_(self, request):
        return 'rId2'

    @pytest.fixture
    def rIds_(self, request, rId_, rId_2_):
        return rId_, rId_2_

    @pytest.fixture
    def Slide_(self, request, slide_):
        Slide_ = class_mock(request, 'pptx.parts.slides.Slide')
        Slide_.new.return_value = slide_
        return Slide_

    @pytest.fixture
    def sldId_(self, request, rId_):
        return instance_mock(request, CT_SlideId, rId=rId_)

    @pytest.fixture
    def sldId_2_(self, request, rId_2_):
        return instance_mock(request, CT_SlideId, rId=rId_2_)

    @pytest.fixture
    def sldIdLst_(self, request, sldId_, sldId_2_):
        sldIdLst_ = instance_mock(request, CT_SlideIdList)
        sldIdLst_.__getitem__.side_effect = [sldId_, sldId_2_]
        sldIdLst_.__iter__.return_value = iter([sldId_, sldId_2_])
        sldIdLst_.__len__.return_value = 2
        return sldIdLst_

    @pytest.fixture
    def slide_(self, request):
        return instance_mock(request, Slide)

    @pytest.fixture
    def slide_2_(self, request):
        return instance_mock(request, Slide)

    @pytest.fixture
    def slide_parts_(self, request, slide_, slide_2_):
        return slide_, slide_2_

    @pytest.fixture
    def slidelayout_(self, request):
        return instance_mock(request, SlideLayout)

    @pytest.fixture
    def slides(self, sldIdLst_, prs_):
        return SlideCollection(sldIdLst_, prs_)

    @pytest.fixture
    def slides_with_slide_parts_(self, sldIdLst_, prs_, slide_parts_):
        slide_, slide_2_ = slide_parts_
        slides = SlideCollection(sldIdLst_, prs_)
        return slides, slide_, slide_2_


class DescribeSlideLayout(object):

    def _loaded_slidelayout(self, prs_slidemaster=None):
        """
        Return SlideLayout instance loaded using mocks. *prs_slidemaster* is
        an already-loaded model-side SlideMaster instance (or mock, as
        appropriate to calling test).
        """
        # partname for related slideMaster
        sldmaster_partname = '/ppt/slideMasters/slideMaster1.xml'
        # path to test slideLayout XML
        slidelayout_path = absjoin(test_file_dir, 'slideLayout1.xml')
        # model-side slideMaster part
        if prs_slidemaster is None:
            prs_slidemaster = Mock(spec=SlideMaster)
        # a part dict containing the already-loaded model-side slideMaster
        loaded_part_dict = {sldmaster_partname: prs_slidemaster}
        # a slideMaster package part for rel target
        pkg_slidemaster_part = Mock(spec=package.Part)
        pkg_slidemaster_part.partname = sldmaster_partname
        # a package-side relationship from slideLayout to its slideMaster
        rel = Mock(name='pptx.package.Relationship')
        rel.rId = 'rId1'
        rel.reltype = RT.SLIDE_MASTER
        rel.target = pkg_slidemaster_part
        # the slideLayout package part to send to _load()
        pkg_slidelayout_part = Mock(spec=package.Part)
        pkg_slidelayout_part.relationships = [rel]
        with open(slidelayout_path, 'rb') as f:
            pkg_slidelayout_part.blob = f.read()
        # _load and return
        slidelayout = SlideLayout()
        return slidelayout._load(pkg_slidelayout_part, loaded_part_dict)

    # def test__load_sets_slidemaster(self):
    #     """SlideLayout._load() sets slidemaster"""
    #     # setup ------------------------
    #     prs_slidemaster = Mock(spec=SlideMaster)
    #     # exercise ---------------------
    #     loaded_slidelayout = self._loaded_slidelayout(prs_slidemaster)
    #     # verify -----------------------
    #     expected = prs_slidemaster
    #     actual = loaded_slidelayout.slidemaster
    #     msg = "expected: %s, got %s" % (expected, actual)
    #     assert actual == expected, msg

    # def test_slidemaster_raises_on_ref_before_assigned(self, slidelayout):
    #     """SlideLayout.slidemaster raises on referenced before assigned"""
    #     with pytest.raises(AssertionError):
    #         slidelayout.slidemaster

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def slidelayout(self):
        return SlideLayout()


class DescribeSlideMaster(object):

    def test_slidelayouts_property_empty_on_construction(self, slidemaster):
        assert len(slidemaster.slidelayouts) == 0

    def test_slidelayouts_correct_length_after_open(self):
        slidemaster = Package.open(test_pptx_path).presentation.slidemasters[0]
        slidelayouts = slidemaster.slidelayouts
        assert len(slidelayouts) == 11

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def slidemaster(self):
        partname = PackURI('/ppt/slideMasters/slideMaster1.xml')
        return SlideMaster(partname, None, None, None)
